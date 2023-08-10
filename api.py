from langchain.document_loaders import Docx2txtLoader
from fastapi import FastAPI, File, UploadFile, Form 
from langchain.document_loaders import PyPDFLoader
from pydantic import BaseModel
from pptx import Presentation
from typing import Annotated
from docx import Document
import tempfile
import prompt
import os
import uvicorn


class Prompt(BaseModel):
    prompt: str

app = FastAPI()

@app.get("/health")
async def health():
    return {"Message": "Up and Running"}

@app.post("/")
async def create_item(res: Prompt):
    return res


@app.post("/chat")
async def upload(query: Annotated[str, Form()], file: Annotated[UploadFile | None, File()] = None):
    if not file:
        response = prompt.generate_response(query)
        return {"response": response}
    else:
        temp_dir = tempfile.mkdtemp()
        file_path = os.path.join(temp_dir, file.filename)

        with open(file_path, "wb") as f:
            f.write(file.file.read())

        if (file.filename.endswith(".pdf")):
            loader = PyPDFLoader(file_path)
            pages = loader.load_and_split()
            response = prompt.generate_response_from_pdf_file(pages, query)


        elif (file.filename.endswith(".docx") or file.filename.endswith(".doc")):
            loader = Docx2txtLoader(file_path)
            data = loader.load()
            response = prompt.generate_response_from_file(data, query)


                
        elif (file.filename.endswith(".ppt") or file.filename.endswith(".pptx")):
            print("It's a ppt file")
            content = ""
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        content +=  shape.text + "\n"
            print(content)  

        elif (file.filename.endswith(".txt")):
            print("It's a txt file")
            content = ""
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read() 
                print(content)

        return {"response": response}
    

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)